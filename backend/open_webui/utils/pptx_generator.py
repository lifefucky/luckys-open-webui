# backend/open_webui/utils/pptx_generator.py
import re
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Any

from markdown import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from open_webui.env import STATIC_DIR
from open_webui.models.chats import ChatTitleMessagesForm


class PPTXGenerator:
    """Generate a PowerPoint (.pptx) presentation from chat messages.
    Mirrors the API of PDFGenerator so it can be used as a drop‑in component.
    """

    def __init__(self, form_data: ChatTitleMessagesForm):
        self.form_data = form_data
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.title_font = Pt(44)
        self.header_font = Pt(32)
        self.body_font = Pt(18)
        self.assets_dir = Path(STATIC_DIR / "assets")

    # ------------------------------------------------------------------
    # Helpers (similar to PDFGenerator)
    # ------------------------------------------------------------------
    def _format_timestamp(self, ts: float) -> str:
        try:
            return datetime.fromtimestamp(ts).strftime("%Y-%m-%d, %H:%M:%S")
        except Exception:
            return ""

    def _markdown_to_text(self, md: str) -> List[str]:
        """Very light markdown → plain‑text conversion returning bullet lines."""
        html = markdown(md)
        for tag in ("<br>", "<br/>", "<p>", "</p>", "<li>", "</li>"):
            html = html.replace(tag, "\n")
        text = re.sub(r"<[^>]+>", "", html)
        return [ln.strip() for ln in text.splitlines() if ln.strip()]

    # ------------------------------------------------------------------
    # Slide creation methods
    # ------------------------------------------------------------------
    def _add_title_slide(self, title: str, subtitle: str = "", date: str = ""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = self.title_font
        if subtitle or date:
            placeholder = slide.placeholders[1]
            txt = subtitle
            if date:
                txt = f"{txt}\n{date}" if txt else date
            placeholder.text = txt
            placeholder.text_frame.paragraphs[0].font.size = Pt(24)
            placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        return slide

    def _add_section_slide(self, section_title: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = section_title
        slide.shapes.title.text_frame.paragraphs[0].font.size = self.header_font
        return slide

    def _add_content_slide(self, title: str, bullets: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = self.header_font
        body = slide.placeholders[1].text_frame
        for b in bullets[:5]:  # max 5 bullets per slide
            p = body.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = self.body_font
        return slide

    # ------------------------------------------------------------------
    # Public interface – mimics PDFGenerator.generate_chat_pdf()
    # ------------------------------------------------------------------
    def generate_chat_pptx(self) -> bytes:
        """Build a PPTX from self.form_data and return the binary content."""
        # Title slide
        self._add_title_slide(self.form_data.title, date=datetime.now().strftime("%Y-%m-%d"))
        # Iterate over messages – each becomes a section + content slide
        for msg in self.form_data.messages:
            role = msg.get("role", "user").title()
            content_md = msg.get("content", "")
            ts = self._format_timestamp(msg.get("timestamp", 0))
            header = f"{role} – {ts}".strip(" –")
            self._add_section_slide(header)
            bullets = self._markdown_to_text(content_md) or [content_md]
            self._add_content_slide("Message", bullets)
        out = BytesIO()
        self.prs.save(out)
        return out.getvalue()
